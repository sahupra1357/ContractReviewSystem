from pptx import Presentation

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return "\n".join(text_content)


if __name__ == "__main__":
    pptx_file_path = "AI_Contract_CoPilot_Security_Review.pptx"  # Replace with your .pptx file path
    extracted_text = extract_text_from_pptx(pptx_file_path)
    with open("AI_Contract_CoPilot_Security_Review.txt", "w", encoding="utf-8") as f:
        f.write(extracted_text)

    print(extracted_text)  # Print the extracted text to the console
